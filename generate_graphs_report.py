import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from PIL import Image
import os

# --- Ler os dados ---
data = pd.read_excel("team_kpis_mock_1.xlsx")  # ajuste para seu arquivo real

# 1. Pie chart: Count KPIs by Status
status_counts = data['Status'].value_counts()

plt.figure(figsize=(6,6))
plt.pie(status_counts, labels=status_counts.index, autopct='%1.1f%%', startangle=140, colors=[ '#FFC107','#4CAF50', '#F44336'])
plt.title('KPIs by Status')
plt.tight_layout()
plt.savefig("kpis_by_status.png")
plt.close()

# 2. Bar chart: Number of KPIs by Responsible
resp_counts = data['Responsible'].value_counts().sort_values(ascending=True)

plt.figure(figsize=(8,5))
resp_counts.plot(kind='barh', color='skyblue')
plt.xlabel('Number of KPIs')
plt.ylabel('Responsible')
plt.title('KPIs by Responsible')
plt.tight_layout()
plt.savefig("kpis_by_responsible.png")
plt.close()

# 3. Bar chart: Number of KPIs by Category
categories_counts = data['Category'].value_counts().sort_values(ascending=True)

plt.figure(figsize=(8,5))
categories_counts.plot(kind='barh', color='lightcoral')
plt.xlabel('Number of KPIs')
plt.ylabel('Category')
plt.title('KPIs by Category')
plt.tight_layout()
plt.savefig("kpis_by_category.png")
plt.close()

charts = [
    "kpis_by_status.png",
    "kpis_by_responsible.png",
    "kpis_by_category.png"
]

# --- 4. Criar PowerPoint com 3 gráficos ---
ppt_path = "KPI_Report.pptx"
prs = Presentation()

for i, chart_path in enumerate(charts):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # slide em branco
    # Adicionar título do slide conforme o gráfico
    title = slide.shapes.title
    if i == 0:
        title.text = "KPIs by Status"
    elif i == 1:
        title.text = "KPIs by Responsible"
    else:
        title.text = "KPIs by Category"

    # Adicionar gráfico na posição e tamanho desejados
    slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(8))
    
prs.save(ppt_path)

# --- 5. Criar PDF com 3 gráficos ---
pdf_path = "KPI_Report.pdf"
c = canvas.Canvas(pdf_path, pagesize=letter)
c.setFont("Helvetica-Bold", 20)
width_page, height_page = letter

for i, chart_path in enumerate(charts):
    if i > 0:
        c.showPage()
    if i == 0:
        title = "KPIs by Status"
    elif i == 1:
        title = "KPIs by Responsible"
    else:
        title = "KPIs by Category"
    
    c.drawString(100, 750, title)
    
    # Pega tamanho original da imagem
    img = Image.open(chart_path)
    img_width, img_height = img.size
    
    max_width = 500
    scale = max_width / img_width
    new_width = max_width
    new_height = img_height * scale
    
    x = 50
    y = height_page - new_height - 100
    
    c.drawImage(chart_path, x, y, width=new_width, height=new_height)

c.save()

print(f"✅ PowerPoint saved as {ppt_path}")
print(f"✅ PDF saved as {pdf_path}")
